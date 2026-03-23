#!/usr/bin/env python3
"""Generate ODH Platform Chaos presentation PPTX from RedHat template."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = "Red Hat OpenShift _ Presentation template.pptx"
OUTPUT = "ODH-Platform-Chaos-Architects-Council.pptx"
IMAGES_DIR = "images"

# Layout indices — verified placeholder positions
LY_TITLE = 0         # TITLE: ph0 title@(2.3,1.9) 8.2x1.8, ph1 subtitle@(2.3,3.9)
LY_CLOSING = 1       # TITLE_1: ph0 title@(2.3,0.6) 5.4x2.2, ph1 subtitle@(2.3,3.3) 4.5x1.7
LY_DIVIDER = 14      # CUSTOM_8: ph0 title@(2.3,2.1) 8.7x3.4, ph1 subtitle@(-2.3,2.3)
LY_IMAGE = 16        # CUSTOM_4: ph0 title@(1.0,1.2) 11.4x0.6, ph3 subtitle@(1.0,1.8)
LY_CONTENT = 21      # CUSTOM_4_17: ph0 title@(1.0,1.2), ph2 subtitle@(1.0,1.8), ph4 body@(2.7,2.5) 8x4
LY_TWOCOL = 24       # CUSTOM_4_17_1: ph0 title@(1.0,1.2), ph2 body_L@(1.0,2.5) 5.4x4.1, ph3 body_R@(7.0,2.5) 5.4x4.1

INCH = 914400


def remove_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def set_ph(slide, idx, text, font_size=None, bold=False, color=None, font_name=None):
    """Set placeholder text."""
    try:
        ph = slide.placeholders[idx]
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        if font_size:
            run.font.size = Pt(font_size)
        if bold:
            run.font.bold = True
        if color:
            run.font.color.rgb = RGBColor(*color)
        if font_name:
            run.font.name = font_name
    except (KeyError, IndexError):
        pass


def set_body(slide, ph_idx, lines, font_size=12, font_name="Red Hat Text"):
    """Set body placeholder with multiple formatted lines. lines = [(text, bold, color), ...]"""
    try:
        ph = slide.placeholders[ph_idx]
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, item in enumerate(lines):
            text = item[0]
            is_bold = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else (30, 30, 30)
            sz = item[3] if len(item) > 3 else font_size

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(2)
            p.space_before = Pt(0)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(sz)
            run.font.bold = is_bold
            run.font.color.rgb = RGBColor(*clr)
            run.font.name = font_name
    except (KeyError, IndexError):
        pass


def set_body_code(slide, ph_idx, code, font_size=10):
    """Set body placeholder with monospace code."""
    lines = [(line,) for line in code.split('\n')]
    set_body(slide, ph_idx, lines, font_size=font_size, font_name="Courier New")


def add_image(slide, image_path, left_in, top_in, width_in, height_in):
    """Add an image to the slide (dimensions in inches)."""
    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            Inches(left_in), Inches(top_in),
            Inches(width_in), Inches(height_in)
        )


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, TEMPLATE)
    output_path = os.path.join(script_dir, OUTPUT)
    images_dir = os.path.join(script_dir, IMAGES_DIR)

    prs = Presentation(template_path)
    remove_all_slides(prs)
    layouts = prs.slide_layouts

    # ===== SLIDE 1: Title (LY_TITLE) =====
    # ph0=title@(2.3,1.9), ph1=subtitle@(2.3,3.9), ph2=presenter@(2.3,5.0), ph3=date@(4.9,5.0)
    s = prs.slides.add_slide(layouts[LY_TITLE])
    set_ph(s, 0, "ODH Platform Chaos", font_size=40, bold=True)
    set_ph(s, 1, "Semantic Chaos Engineering for OpenDataHub Operators", font_size=20)
    set_ph(s, 2, "Architects Council", font_size=14)
    set_ph(s, 3, "March 2026", font_size=14)

    # ===== SLIDE 2: The Problem (LY_CONTENT) =====
    # ph0=title@(1.0,1.2) 11.4x0.6, ph2=subtitle@(1.0,1.8), ph4=body@(2.7,2.5) 8x4
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "The Problem We Solve", font_size=28, bold=True)
    set_ph(s, 2, "Moving beyond \"does the pod restart?\"", font_size=14)
    set_body(s, 4, [
        ("Traditional chaos tools answer:", False),
        ("\"Does the pod come back after being killed?\"", True, (100, 100, 100)),
        ("", False),
        ("We answer a harder question:", False),
        ("\"After a fault, does the operator semantically restore all managed", True, (200, 0, 0)),
        ("resources — with correct specs, owner references, labels, and", True, (200, 0, 0)),
        ("reconciliation state?\"", True, (200, 0, 0)),
        ("", False),
        ("A pod restarting is necessary but not sufficient. An operator can", False),
        ("restart and still leave ConfigMaps drifted, webhooks misconfigured,", False),
        ("or RBAC bindings empty.", False),
    ], font_size=13)

    # ===== SLIDE 3: Why This Matters (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "Why This Matters for RHOAI", font_size=28, bold=True)
    set_ph(s, 2, "RHOAI operators manage complex resource graphs", font_size=14)
    set_body(s, 4, [
        ("odh-model-controller:  6 managed resources  |  7 webhooks  |  3 finalizers", True),
        ("kserve:  16 managed resources  |  12 webhooks  |  4 finalizers", True),
        ("", False),
        ("(Counts from operator knowledge models — verified against deployed manifests)", False, (120, 120, 120), 11),
        ("", False),
        ("A single unreconciled ConfigMap (inferenceservice-config) can silently", True, (200, 0, 0)),
        ("break ALL InferenceService deployments across the cluster.", True, (200, 0, 0)),
        ("No pod restart. No alerts.", True, (200, 0, 0)),
        ("", False),
        ("This tool catches those failures before users do.", False),
    ], font_size=14)

    # ===== SLIDE 4: Who Uses This (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "Who Uses This and When", font_size=28, bold=True)
    set_ph(s, 2, "Four personas, four usage modes", font_size=14)
    set_body(s, 4, [
        ("Operator developer", True),
        ("  SDK middleware in integration tests — every PR", False),
        ("QE engineer", True),
        ("  CLI suite against staging — release gating", False),
        ("SRE / on-call", True),
        ("  CLI run against production-like env — incident prep", False),
        ("CI pipeline", True),
        ("  Container image + JUnit reports — automated promotion gating", False),
        ("", False),
        ("Common thread: verify semantic healing, not just pod restarts.", True, (200, 0, 0)),
    ], font_size=13)

    # ===== SLIDE 5: Why Build Not Extend (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "Why Build, Not Extend?", font_size=28, bold=True)
    set_ph(s, 2, "Complementary, not competing", font_size=14)
    set_body(s, 4, [
        ("What neither LitmusChaos nor Chaos Mesh does:", True),
        ("  \u2713  Semantic reconciliation verification", False),
        ("  \u2713  Knowledge model (operator-specific oracle)", False),
        ("  \u2713  SDK middleware for integration tests", False),
        ("  \u2713  Fuzz testing (no cluster needed)", False),
        ("  \u2713  Lightweight (no CRD operator to install)", False),
        ("", False),
        ("What all three do:  \u2713  Pod/network fault injection", False),
        ("", False),
        ("They test the platform. We test the operator.", True, (200, 0, 0)),
        ("The two categories are complementary.", False),
    ], font_size=12)

    # ===== DIVIDER: Architecture =====
    s = prs.slides.add_slide(layouts[LY_DIVIDER])
    set_ph(s, 0, "Architecture\nHow it works under the hood", font_size=32, bold=True)

    # ===== SLIDE 6: Architecture Overview (LY_IMAGE) =====
    # ph0=title@(1.0,1.2) 11.4x0.6, ph3=subtitle@(1.0,1.8) 11.4x0.3
    s = prs.slides.add_slide(layouts[LY_IMAGE])
    set_ph(s, 0, "Architecture Overview", font_size=28, bold=True)
    set_ph(s, 3, "All components injected via OrchestratorConfig — no globals, no singletons", font_size=11)
    add_image(s, os.path.join(images_dir, "Architecture Overview.png"), 1.5, 2.3, 10.0, 4.2)

    # ===== SLIDE 7: Experiment Lifecycle (LY_IMAGE) =====
    s = prs.slides.add_slide(layouts[LY_IMAGE])
    set_ph(s, 0, "Experiment Lifecycle — State Machine", font_size=28, bold=True)
    set_ph(s, 3, "Cleanup runs in defer — on cancellation, timeout, or panic. SIGKILL/OOM: annotations survive for clean command.", font_size=11)
    add_image(s, os.path.join(images_dir, "Experiment Lifecycle.png"), 1.0, 2.3, 11.0, 3.8)

    # ===== SLIDE 8: Injection Registry (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "Injection Registry — Strategy Pattern", font_size=28, bold=True)
    set_ph(s, 2, "Every injector returns a CleanupFunc that reverses the fault", font_size=14)
    set_body_code(s, 4,
        'type Injector interface {\n'
        '    Validate(spec InjectionSpec, blast BlastRadiusSpec) error\n'
        '    Inject(ctx context.Context, spec InjectionSpec, ns string)\n'
        '           (CleanupFunc, []InjectionEvent, error)\n'
        '}\n'
        '\n'
        'type CleanupFunc func(ctx context.Context) error\n'
        '\n'
        '// Orchestrator calls cleanup in defer,\n'
        '// guaranteeing rollback even on failure.\n'
        '\n'
        '// Why an interface? Each injection type has different\n'
        '// K8s API interactions (pods vs RBAC vs webhooks).\n'
        '// Common interface enables the registry pattern.',
        font_size=11)

    # ===== SLIDE 9: Seven Injection Types (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "The Seven Injection Types", font_size=28, bold=True)
    set_ph(s, 2, "Five operator failure categories covered", font_size=14)
    set_body_code(s, 4,
        'Type               What It Does                        Danger\n'
        '─────────────────────────────────────────────────────────────\n'
        'PodKill            Force-delete pods (0s grace)         Low\n'
        'NetworkPartition   Deny-all NetworkPolicy               Medium\n'
        'ConfigDrift        Modify ConfigMap/Secret data          Medium\n'
        'CRDMutation        Patch custom resource fields          Medium\n'
        'WebhookDisrupt     Change webhook failure policy         High\n'
        'RBACRevoke         Clear binding subjects                High\n'
        'FinalizerBlock     Add blocking finalizer                Medium\n'
        '\n'
        'Categories: Pod lifecycle | Network | Configuration\n'
        '            Control plane | Lifecycle hooks',
        font_size=11)

    # ===== SLIDE 10: Safety Architecture (LY_IMAGE) =====
    s = prs.slides.add_slide(layouts[LY_IMAGE])
    set_ph(s, 0, "Safety Architecture — Defense in Depth", font_size=28, bold=True)
    set_ph(s, 3, "Every mutation must be reversible. Every artifact must be traceable.", font_size=11)
    add_image(s, os.path.join(images_dir, "Safety Architecture.png"), 2.0, 2.2, 9.0, 4.3)

    # ===== SLIDE 11: Security & Audit (LY_TWOCOL) =====
    # ph0=title@(1.0,1.2), ph2=bodyL@(1.0,2.5) 5.4x4.1, ph3=bodyR@(7.0,2.5) 5.4x4.1
    s = prs.slides.add_slide(layouts[LY_TWOCOL])
    set_ph(s, 0, "Security, Authorization & Audit", font_size=28, bold=True)
    set_ph(s, 5, "RBAC is the security boundary — no custom auth layer", font_size=11)
    set_body(s, 2, [
        ("Required permissions:", True),
        ("", False),
        ("PodKill: delete on pods", False),
        ("ConfigDrift: get, patch on configmaps/secrets", False),
        ("WebhookDisrupt: get, patch on webhook configs", False),
        ("RBACRevoke: get, patch on clusterrolebindings", False),
        ("DistributedLock: create, get, update on leases", False),
        ("Reports: create, update on configmaps", False),
        ("", False),
        ("Why no custom auth layer?", True),
        ("K8s RBAC already solves authorization.", False),
    ], font_size=12)
    set_body(s, 3, [
        ("Audit trail:", True),
        ("", False),
        ("Every experiment produces a report ConfigMap:", False),
        ("  \u2022 Experiment name + timestamp", False),
        ("  \u2022 Verdict + injection details", False),
        ("  \u2022 Recovery metrics", False),
        ("", False),
        ("Labeled + queryable:", True),
        ("chaos.opendatahub.io/verdict={v}", False, (30, 30, 30), 10),
        ("managed-by=odh-chaos", False, (30, 30, 30), 10),
        ("", False),
        ("K8s audit logs capture which SA created them.", False),
    ], font_size=12)

    # ===== DIVIDER: Knowledge Model =====
    s = prs.slides.add_slide(layouts[LY_DIVIDER])
    set_ph(s, 0, "Knowledge Model\nTeaching the system what operators should reconcile", font_size=32, bold=True)

    # ===== SLIDE 12: Knowledge Model YAML (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "Knowledge Model — Encoding Operator Semantics", font_size=26, bold=True)
    set_ph(s, 2, "30-80 lines of YAML per operator", font_size=14)
    set_body_code(s, 4,
        'operator:\n'
        '  name: odh-model-controller\n'
        '  namespace: opendatahub\n'
        'components:\n'
        '  - name: odh-model-controller\n'
        '    controller: DataScienceCluster\n'
        '    managedResources:\n'
        '      - kind: Deployment\n'
        '        name: odh-model-controller\n'
        '        expectedSpec: { replicas: 1 }\n'
        '    webhooks:\n'
        '      - name: validating.isvc.odh-model-controller\n'
        '    finalizers: [odh.inferenceservice.finalizers]\n'
        '    steadyState:\n'
        '      checks:\n'
        '        - type: conditionTrue\n'
        '          conditionType: Available\n'
        'recovery:\n'
        '  reconcileTimeout: "300s"\n'
        '  maxReconcileCycles: 10',
        font_size=11)

    # ===== SLIDE 13: What KM Enables (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "What the Knowledge Model Enables", font_size=28, bold=True)
    set_ph(s, 2, "From \"did it crash?\" to \"did it heal correctly?\"", font_size=14)
    set_body(s, 4, [
        ("Without it: kill a pod, check if it restarts.", False),
        ("", False),
        ("With it, we verify:", True),
        ("  \u2022 Did the operator reconcile all 6 managed resources?", False),
        ("  \u2022 Are owner references intact?", False),
        ("  \u2022 Is expectedSpec.replicas still correct?", False),
        ("  \u2022 Did the webhook get re-registered?", False),
        ("  \u2022 Was recovery within 300s timeout and 10 cycle limit?", False),
        ("", False),
        ("Maintenance: 30-80 lines YAML, versioned with operator code.", False),
        ("Updates needed roughly once per release cycle.", False),
        ("", False),
        ("OLM: references operator-created resources (not CSV resources).", False, (120, 120, 120)),
    ], font_size=13)

    # ===== SLIDE 14: Verdict Engine (LY_IMAGE) =====
    s = prs.slides.add_slide(layouts[LY_IMAGE])
    set_ph(s, 0, "Verdict Engine — Decision Tree", font_size=28, bold=True)
    set_ph(s, 3, "Four verdicts, not two. Degraded = \"recovered, but not well enough\" — actionable without blocking releases.", font_size=11)
    add_image(s, os.path.join(images_dir, "Verdict Engine.png"), 2.65, 2.2, 8.0, 4.3)

    # ===== SLIDE 15: Three Usage Modes (LY_IMAGE) =====
    s = prs.slides.add_slide(layouts[LY_IMAGE])
    set_ph(s, 0, "Three Usage Modes", font_size=28, bold=True)
    set_ph(s, 3, "One knowledge model drives all three modes.", font_size=11)
    add_image(s, os.path.join(images_dir, "Three Usage Modes.png"), 1.0, 2.3, 11.0, 4.3)

    # ===== SLIDE 16: SDK Middleware (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "SDK Middleware — Testing Without a Cluster", font_size=26, bold=True)
    set_ph(s, 2, "ChaosClient implements client.Client — drop-in replacement", font_size=14)
    set_body_code(s, 4,
        '// Wrap any controller-runtime client with fault injection\n'
        'chaosClient := sdk.NewChaosClient(realClient, faultConfig)\n'
        '\n'
        '// Every K8s API call goes through MaybeInject()\n'
        'func (c *ChaosClient) Get(ctx, key, obj, opts...) error {\n'
        '    if err := c.faults.MaybeInject(OpGet); err != nil {\n'
        '        return err  // Injected fault\n'
        '    }\n'
        '    return c.inner.Get(ctx, key, obj, opts...)\n'
        '}\n'
        '\n'
        '// Performance: single map lookup + random check per call\n'
        '// Why not mocks? Middleware injects faults probabilistically\n'
        '// into real code paths, testing actual error handling.',
        font_size=12)

    # ===== DIVIDER: In Practice =====
    s = prs.slides.add_slide(layouts[LY_DIVIDER])
    set_ph(s, 0, "In Practice\nReal experiments, real findings", font_size=32, bold=True)

    # ===== SLIDE 17: Concrete Example (LY_TWOCOL) =====
    s = prs.slides.add_slide(layouts[LY_TWOCOL])
    set_ph(s, 0, "ConfigDrift on inferenceservice-config", font_size=24, bold=True)
    set_ph(s, 5, "A real experiment with a real finding", font_size=11)
    set_body_code(s, 2,
        'metadata:\n'
        '  name: config-drift-isvc\n'
        'spec:\n'
        '  injection:\n'
        '    type: ConfigDrift\n'
        '    target:\n'
        '      name: inferenceservice-config\n'
        '      namespace: opendatahub\n'
        '    parameters:\n'
        '      key: "config"\n'
        '      value: \'{"corrupted": true}\'\n'
        '  blastRadius:\n'
        '    allowedNamespaces:\n'
        '      - opendatahub\n'
        '    maxPodsAffected: 0',
        font_size=11)
    set_body(s, 3, [
        ("What we found:", True, (200, 0, 0)),
        ("", False),
        ("The controller did NOT", True, (200, 0, 0)),
        ("reconcile this ConfigMap.", True, (200, 0, 0)),
        ("", False),
        ("It is created during", False),
        ("installation but never", False),
        ("watched for drift.", False),
        ("", False),
        ("Any manual edit silently", False),
        ("breaks every InferenceService", False),
        ("deployment on the cluster.", False),
        ("", False),
        ("This finding drove a fix", True),
        ("upstream.", True),
    ], font_size=12)

    # ===== SLIDE 18: Tool Failure Modes (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "Tool Failure Modes — What If We Fail?", font_size=26, bold=True)
    set_ph(s, 2, "The tool must never make things worse than the fault it injected", font_size=14)
    set_body(s, 4, [
        ("Tool crashes mid-injection", True),
        ("  \u2192 defer cleanup runs. SIGKILL/OOM: rollback annotations persist, clean restores.", False),
        ("", False),
        ("Rollback annotation corrupted", True),
        ("  \u2192 SHA-256 checksum refuses to apply, logs for manual recovery.", False),
        ("", False),
        ("Lock lease not released", True),
        ("  \u2192 15-minute auto-expiry via leaseDurationSeconds.", False),
        ("", False),
        ("Tool loses cluster connectivity", True),
        ("  \u2192 Rollback annotations persist. TTL marks artifacts for GC.", False),
        ("", False),
        ("Knowledge model is wrong", True),
        ("  \u2192 preflight validates against live cluster before any experiment.", False),
    ], font_size=12)

    # ===== SLIDE 19: CI Integration (LY_TWOCOL) =====
    s = prs.slides.add_slide(layouts[LY_TWOCOL])
    set_ph(s, 0, "CI Integration & Reports", font_size=28, bold=True)
    set_ph(s, 5, "Designed from day one as a CI citizen", font_size=11)
    set_body(s, 2, [
        ("Exit code contract:", True),
        ("", False),
        ("preflight  \u2192 Exit 0 = resources OK", False),
        ("run        \u2192 Exit 0 = Resilient", False),
        ("suite      \u2192 Exit 0 = all pass", False),
        ("validate   \u2192 Exit 0 = valid YAML", False),
        ("", False),
        ("Report formats:", True),
        ("  \u2022 JSON files", False),
        ("  \u2022 JUnit XML", False),
        ("  \u2022 K8s ConfigMaps", False),
    ], font_size=12)
    set_body(s, 3, [
        ("Container image:", True),
        ("", False),
        ("  distroless/static:nonroot", False),
        ("  Non-root UID 65532", False),
        ("  Multi-arch amd64/arm64", False),
        ("", False),
        ("CI docs provided:", True),
        ("  \u2022 Tekton Tasks + Pipelines", False),
        ("  \u2022 GitHub Actions", False),
        ("  \u2022 Deployment gating", False),
    ], font_size=12)

    # ===== SLIDE 20: Non-Goals (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "What We Don't Do (Deliberate Non-Goals)", font_size=26, bold=True)
    set_ph(s, 2, "Scope awareness", font_size=14)
    set_body(s, 4, [
        ("CRD-based controller", True),
        ("  CLI-first avoids premature API commitment. CRD mode planned after stabilization.", False),
        ("", False),
        ("Real-time dashboard", True),
        ("  Reports via ConfigMaps + JUnit. UI contradicts \"run and exit\" model.", False),
        ("", False),
        ("Network-level injection (iptables/tc)", True),
        ("  NetworkPolicy is K8s-native. iptables requires privileged containers.", False),
        ("", False),
        ("Multi-cluster support", True),
        ("  Single-cluster keeps safety simple. Cross-cluster = different problem.", False),
        ("", False),
        ("Mutation webhooks for injection", True),
        ("  Too invasive — blast radius unbounded. Direct API mutations are reversible.", False),
    ], font_size=12)

    # ===== DIVIDER: Roadmap =====
    s = prs.slides.add_slide(layouts[LY_DIVIDER])
    set_ph(s, 0, "Roadmap & Closing\nWhere we go from here", font_size=32, bold=True)

    # ===== SLIDE 21: Roadmap & Ask (LY_TWOCOL) =====
    s = prs.slides.add_slide(layouts[LY_TWOCOL])
    set_ph(s, 0, "Roadmap & Ask", font_size=28, bold=True)
    set_ph(s, 5, "What's done, what's next, what we need from you", font_size=11)
    set_body(s, 2, [
        ("Done:", True, (34, 139, 34)),
        ("\u2713 Core framework", False),
        ("  7 injectors, orchestrator,", False),
        ("  evaluator, safety", False),
        ("\u2713 SDK middleware", False),
        ("  ChaosClient, WrapReconciler", False),
        ("\u2713 CLI & CI", False),
        ("  10 commands, JUnit, container", False),
        ("", False),
        ("Planned:", True, (200, 0, 0)),
        ("\u2192 Knowledge models for all", False),
        ("   RHOAI operators", False),
        ("\u2192 ChaosExperiment CRD", False),
        ("\u2192 OpenShift CI integration", False),
    ], font_size=12)
    set_body(s, 3, [
        ("Our ask:", True, (200, 0, 0)),
        ("", False),
        ("1. Adopt as recommended", True),
        ("   practice", True),
        ("   Pilot 2-3 operators this", False),
        ("   quarter, evaluate results", False),
        ("", False),
        ("2. Advisory-mode chaos", True),
        ("   suites in CI", True),
        ("   Informational for 2 cycles", False),
        ("   before gating", False),
        ("", False),
        ("3. Feedback on operator", True),
        ("   coverage priorities", True),
    ], font_size=12)

    # ===== SLIDE 22: Summary (LY_CONTENT) =====
    s = prs.slides.add_slide(layouts[LY_CONTENT])
    set_ph(s, 0, "Summary", font_size=28, bold=True)
    set_ph(s, 2, "ODH Platform Chaos tests what matters: semantic reconciliation correctness", font_size=14)
    set_body(s, 4, [
        ("Key architectural decisions:", True),
        ("", False),
        ("1. Knowledge-driven — operator semantics encoded in YAML", False),
        ("2. Safety-first — 6 layers of defense in depth", False),
        ("3. Interface-driven — Injector, Observer, Lock are pluggable", False),
        ("4. Three modes — CLI (cluster), SDK (integration), Fuzz (unit)", False),
        ("5. CI-native — exit codes, JUnit, container image", False),
        ("6. RBAC-delegated — Kubernetes is the security boundary", False),
        ("", False),
        ("One question to take away:", True, (200, 0, 0)),
        ("For each RHOAI operator, can we prove that after any of these", True, (200, 0, 0)),
        ("7 fault types, the operator restores all managed resources to", True, (200, 0, 0)),
        ("their correct state within the expected time?", True, (200, 0, 0)),
    ], font_size=14)

    # ===== SLIDE 23: Thank You (LY_CLOSING) =====
    # ph0=title@(2.3,0.6) 5.4x2.2, ph1=subtitle@(2.3,3.3) 4.5x1.7
    s = prs.slides.add_slide(layouts[LY_CLOSING])
    set_ph(s, 0, "Thank you", font_size=40, bold=True)
    set_ph(s, 1, "Questions?\n\nodh-platform-chaos\ngithub.com/opendatahub-io/odh-platform-chaos", font_size=14)

    # Save
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
